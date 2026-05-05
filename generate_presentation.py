from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0A, 0x29, 0x5C)   # USC navy
GOLD        = RGBColor(0xFF, 0xC7, 0x2C)   # USC gold
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_BLUE    = RGBColor(0x1A, 0x52, 0x96)
GREEN       = RGBColor(0x1B, 0x8A, 0x5A)
RED         = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_BLUE = RGBColor(0x21, 0x96, 0xF3)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_bullet_slide(prs, title_text, bullets, subtitle=None):
    slide = prs.slides.add_slide(BLANK)

    # header bar
    add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
    add_text(slide, title_text, Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # gold accent line
    add_rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)

    # white body
    add_rect(slide, 0, Inches(1.16), W, H - Inches(1.16), LIGHT_GRAY)

    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.45),
                 font_size=16, italic=True, color=MED_BLUE, align=PP_ALIGN.LEFT)

    top = Inches(1.25) + (Inches(0.45) if subtitle else 0)
    for bullet in bullets:
        indent = bullet.startswith("  ")
        btext  = bullet.lstrip()
        fs     = 17 if not indent else 15
        col    = DARK_GRAY if not indent else MED_BLUE
        left   = Inches(0.5) if not indent else Inches(0.85)
        prefix = "▸  " if not indent else "–  "
        add_text(slide, prefix + btext, left, top, Inches(12.0), Inches(0.55),
                 font_size=fs, color=col, align=PP_ALIGN.LEFT)
        top += Inches(0.56) if not indent else Inches(0.5)

    # slide number / footer
    add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK_BLUE)
    return slide


def add_table_slide(prs, title_text, headers, rows, subtitle=None, col_widths=None, highlight_row=None):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
    add_text(slide, title_text, Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85),
             font_size=28, bold=True, color=WHITE)
    add_rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)
    add_rect(slide, 0, Inches(1.16), W, H - Inches(1.16), LIGHT_GRAY)

    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4),
                 font_size=15, italic=True, color=MED_BLUE)

    top = Inches(1.7) if subtitle else Inches(1.3)
    n_cols = len(headers)
    cw = col_widths or [Inches(13.0) / n_cols] * n_cols
    left_start = Inches(0.2)

    # header row
    x = left_start
    for i, h in enumerate(headers):
        add_rect(slide, x, top, cw[i], Inches(0.42), DARK_BLUE)
        add_text(slide, h, x + Inches(0.05), top + Inches(0.03), cw[i] - Inches(0.1), Inches(0.38),
                 font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        x += cw[i]
    top += Inches(0.42)

    for ri, row in enumerate(rows):
        bg = RGBColor(0xE8, 0xF4, 0xFD) if highlight_row == ri else (
             WHITE if ri % 2 == 0 else RGBColor(0xEE, 0xEE, 0xEE))
        x = left_start
        for ci, cell in enumerate(row):
            add_rect(slide, x, top, cw[ci], Inches(0.42), bg,
                     line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
            is_bold = highlight_row == ri
            cell_color = GREEN if highlight_row == ri else DARK_GRAY
            add_text(slide, cell, x + Inches(0.05), top + Inches(0.03),
                     cw[ci] - Inches(0.1), Inches(0.38),
                     font_size=12, bold=is_bold, color=cell_color, align=PP_ALIGN.CENTER)
            x += cw[ci]
        top += Inches(0.42)

    add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK_BLUE)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, DARK_BLUE)
add_rect(slide, 0, Inches(2.8), W, Inches(0.08), GOLD)
add_rect(slide, 0, Inches(5.2), W, Inches(0.08), GOLD)

add_text(slide,
         "Adaptive Fraud Detection",
         Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.0),
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide,
         "Dynamic Thresholding · Cost-Sensitive Learning · OOD Detection",
         Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.7),
         font_size=22, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

add_text(slide,
         "EE559 / CSCI 559 — Machine Learning, Spring 2026\n"
         "Prof. Anand A. Joshi  |  University of Southern California",
         Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.9),
         font_size=17, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)

add_text(slide,
         "Om Suresh Prajapati  |  USC ID: 1430823821",
         Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.6),
         font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide,
         "Final Presentation  ·  May 4, 2026",
         Inches(0.5), Inches(4.75), Inches(12.3), Inches(0.5),
         font_size=16, color=GOLD, align=PP_ALIGN.CENTER)

# decorative corner accent
add_rect(slide, 0, H - Inches(0.4), Inches(3), Inches(0.4), MED_BLUE)
add_rect(slide, W - Inches(3), H - Inches(0.4), Inches(3), Inches(0.4), MED_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide(prs, "Agenda",
    [
        "Problem Statement & Motivation",
        "Dataset Overview",
        "ML Pipeline & Preprocessing",
        "Class Imbalance Strategies",
        "Models & Evaluation Metrics",
        "Results: Baseline vs. Full Comparison",
        "Threshold Tuning",
        "Cost-Sensitive Evaluation",
        "Novel Contribution: OOD Detection (Isolation Forest)",
        "Error Analysis",
        "Final Test Results & Conclusions",
    ],
    subtitle="15-minute overview")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide(prs, "Problem Statement & Motivation",
    [
        "Credit card fraud costs the global economy over $30 billion annually",
        "Fraudulent transactions are extremely rare → severe class imbalance",
        "Standard accuracy is misleading: 99.83% by predicting all legitimate",
        "Asymmetric costs: missing fraud (FN) is far more costly than a false alarm (FP)",
        "  FN cost = €500 per missed fraud   |   FP cost = €10 per false alarm",
        "Goal: maximize fraud recall while controlling false positives under a real cost model",
        "Novel extension: detect previously unseen fraud patterns with OOD detection",
    ],
    subtitle="Why fraud detection is a hard ML problem")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset Overview
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text(slide, "Dataset Overview", Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85),
         font_size=28, bold=True, color=WHITE)
add_rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)
add_rect(slide, 0, Inches(1.16), W, H - Inches(1.16), LIGHT_GRAY)

# Left column — stats boxes
stats = [
    ("284,807", "Total Transactions"),
    ("492", "Fraudulent (0.17%)"),
    ("30", "Features (V1–V28 + Time + Amount)"),
    ("PCA", "V1–V28 are anonymised principal components"),
]
top = Inches(1.4)
for val, label in stats:
    add_rect(slide, Inches(0.3), top, Inches(4.8), Inches(0.85), DARK_BLUE)
    add_text(slide, val, Inches(0.35), top + Inches(0.02), Inches(1.5), Inches(0.55),
             font_size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, label, Inches(1.9), top + Inches(0.15), Inches(3.1), Inches(0.55),
             font_size=14, color=WHITE, align=PP_ALIGN.LEFT)
    top += Inches(1.0)

# Right column — key characteristics
add_rect(slide, Inches(5.6), Inches(1.4), Inches(7.4), Inches(5.6), WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(slide, "Key Characteristics", Inches(5.8), Inches(1.5), Inches(7.0), Inches(0.45),
         font_size=17, bold=True, color=DARK_BLUE)

chars = [
    "Source: Kaggle Credit Card Fraud Detection dataset",
    "Split: 70% train / 10% validation / 20% test",
    "Stratified split preserves 0.17% fraud rate in each split",
    "Amount feature log-transformed (log1p) before scaling",
    "Time feature scaled with StandardScaler",
    "No missing values; no text features",
    "Class ratio  →  1 fraud per ~578 legitimate transactions",
]
t = Inches(2.05)
for c in chars:
    add_text(slide, "•  " + c, Inches(5.8), t, Inches(7.0), Inches(0.48),
             font_size=13, color=DARK_GRAY)
    t += Inches(0.5)

add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ML Pipeline
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text(slide, "ML Pipeline & Preprocessing", Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85),
         font_size=28, bold=True, color=WHITE)
add_rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)
add_rect(slide, 0, Inches(1.16), W, H - Inches(1.16), LIGHT_GRAY)

steps = [
    ("1", "Load Data", "284,807 rows · 30 features"),
    ("2", "Stratified Split", "70 / 10 / 20  train/val/test"),
    ("3", "Feature Engineering", "log1p(Amount) · StandardScaler"),
    ("4", "Imbalance Handling", "Balanced weights / SMOTE / Under-sampling"),
    ("5", "Model Training", "LR · DT · RF · GBT"),
    ("6", "Threshold Tuning", "Sweep on validation set"),
    ("7", "Cost-Sensitive Eval", "€500 FN / €10 FP business cost"),
    ("8", "OOD Detection", "Isolation Forest on legitimate transactions"),
]

box_w = Inches(1.4)
box_h = Inches(0.7)
gap   = Inches(0.15)
start_x = Inches(0.25)
start_y = Inches(1.5)

for i, (num, title, desc) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    # box
    add_rect(slide, x, start_y, box_w, box_h, DARK_BLUE)
    add_text(slide, num, x, start_y, box_w, Inches(0.32),
             font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, title, x, start_y + Inches(0.3), box_w, Inches(0.38),
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # arrow
    if i < len(steps) - 1:
        ax = x + box_w + Inches(0.01)
        add_text(slide, "▶", ax, start_y + Inches(0.15), gap + Inches(0.05), Inches(0.4),
                 font_size=12, color=MED_BLUE, align=PP_ALIGN.CENTER)
    # description below
    add_text(slide, desc, x, start_y + box_h + Inches(0.1), box_w, Inches(0.7),
             font_size=9, color=DARK_GRAY, align=PP_ALIGN.CENTER, italic=True)

# preprocessing details
add_rect(slide, Inches(0.3), Inches(3.5), Inches(12.7), Inches(2.9), WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(slide, "Preprocessing Details", Inches(0.5), Inches(3.6), Inches(6.0), Inches(0.4),
         font_size=16, bold=True, color=DARK_BLUE)
details_l = [
    "Train set only used to fit scaler (prevent data leakage)",
    "SMOTE applied only to training set after split",
    "Validation set used for threshold search; test set touched once",
]
details_r = [
    "PR-AUC preferred over ROC-AUC for imbalanced data",
    "Threshold chosen to maximise F1 or minimise business cost",
    "All models evaluated with same val/test splits for fairness",
]
t = Inches(4.1)
for d in details_l:
    add_text(slide, "•  " + d, Inches(0.5), t, Inches(6.1), Inches(0.42), font_size=12, color=DARK_GRAY)
    t += Inches(0.44)
t = Inches(4.1)
for d in details_r:
    add_text(slide, "•  " + d, Inches(6.8), t, Inches(6.1), Inches(0.42), font_size=12, color=DARK_GRAY)
    t += Inches(0.44)

add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Class Imbalance Strategies
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text(slide, "Class Imbalance Strategies", Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85),
         font_size=28, bold=True, color=WHITE)
add_rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)
add_rect(slide, 0, Inches(1.16), W, H - Inches(1.16), LIGHT_GRAY)

strategies = [
    ("Balanced Class Weights", DARK_BLUE,
     ["Assigns higher penalty to minority (fraud) class during training",
      "Weight ∝ n_samples / (n_classes × n_samples_per_class)",
      "No change to dataset size — computationally cheap",
      "Works natively with sklearn estimators",
      "★ Best performer overall"]),
    ("SMOTE", MED_BLUE,
     ["Synthetic Minority Over-sampling Technique",
      "Generates synthetic fraud samples by interpolating k-NN neighbours",
      "Balances training set to ~50/50 ratio",
      "Risk: synthetic samples may not match real fraud distribution",
      "Best for GBT in our experiments"]),
    ("Random Under-sampling", GREEN,
     ["Randomly drops legitimate transactions from training set",
      "Creates a balanced training set (small and fast)",
      "Loses potentially useful legitimate data",
      "Lower absolute performance — used as a baseline comparison",
      "Useful when training time is constrained"]),
]

col_w = Inches(4.2)
col_h = Inches(5.4)
start_x = Inches(0.2)
for i, (title, color, points) in enumerate(strategies):
    x = start_x + i * (col_w + Inches(0.15))
    add_rect(slide, x, Inches(1.3), col_w, col_h, color)
    add_text(slide, title, x + Inches(0.1), Inches(1.4), col_w - Inches(0.2), Inches(0.5),
             font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    t = Inches(2.05)
    for p in points:
        add_text(slide, "• " + p, x + Inches(0.12), t, col_w - Inches(0.25), Inches(0.65),
                 font_size=11, color=WHITE, align=PP_ALIGN.LEFT)
        t += Inches(0.72)

add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Models & Metrics
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text(slide, "Models & Evaluation Metrics", Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85),
         font_size=28, bold=True, color=WHITE)
add_rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)
add_rect(slide, 0, Inches(1.16), W, H - Inches(1.16), LIGHT_GRAY)

# Models column
add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.0), Inches(5.8), WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(slide, "Models Trained", Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.4),
         font_size=17, bold=True, color=DARK_BLUE)

models = [
    ("Logistic Regression", "Linear baseline; fast, interpretable"),
    ("Decision Tree", "Non-linear, fully interpretable; prone to overfit"),
    ("Random Forest", "Ensemble of trees; robust, high PR-AUC"),
    ("Gradient Boosted Trees", "Sequential boosting; strong on SMOTE"),
    ("Isolation Forest", "Unsupervised OOD detector (novel contribution)"),
]
t = Inches(2.0)
for name, desc in models:
    add_rect(slide, Inches(0.5), t, Inches(5.6), Inches(0.72), DARK_BLUE)
    add_text(slide, name, Inches(0.6), t + Inches(0.04), Inches(5.3), Inches(0.32),
             font_size=13, bold=True, color=GOLD)
    add_text(slide, desc, Inches(0.6), t + Inches(0.35), Inches(5.3), Inches(0.32),
             font_size=11, color=WHITE)
    t += Inches(0.82)

# Metrics column
add_rect(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.8), WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(slide, "Evaluation Metrics", Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.4),
         font_size=17, bold=True, color=DARK_BLUE)

metrics = [
    ("PR-AUC", "PRIMARY", "Area under Precision-Recall curve;\nbest for imbalanced data"),
    ("ROC-AUC", "secondary", "Discrimination ability across\nall thresholds"),
    ("F1-score", "secondary", "Harmonic mean of precision & recall\nat optimal threshold"),
    ("Recall (Sensitivity)", "critical", "Fraction of fraud caught —\nbusiness priority #1"),
    ("Business Cost", "novel", "€500 × FN + €10 × FP\nReal-world cost model"),
]
t = Inches(2.0)
for name, tag, desc in metrics:
    tag_color = RED if tag == "PRIMARY" else (GOLD if tag == "novel" else MED_BLUE)
    add_rect(slide, Inches(7.0), t, Inches(5.6), Inches(0.78), LIGHT_GRAY,
             line_color=RGBColor(0xCC, 0xCC, 0xCC))
    add_rect(slide, Inches(7.0), t, Inches(0.85), Inches(0.78), tag_color)
    add_text(slide, tag, Inches(7.0), t + Inches(0.22), Inches(0.87), Inches(0.35),
             font_size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, name, Inches(7.95), t + Inches(0.04), Inches(4.5), Inches(0.3),
             font_size=13, bold=True, color=DARK_BLUE)
    add_text(slide, desc, Inches(7.95), t + Inches(0.34), Inches(4.6), Inches(0.4),
             font_size=10, color=DARK_GRAY)
    t += Inches(0.88)

add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Baseline Results
# ═══════════════════════════════════════════════════════════════════════════
add_table_slide(prs,
    "Baseline Results (Default Threshold = 0.5)",
    ["Model", "ROC-AUC", "PR-AUC", "Recall", "Precision", "F1-Score"],
    [
        ["Logistic Regression", "0.986", "0.752", "91.9%", "6.8%", "0.127"],
        ["Random Forest",       "0.951", "0.864", "79.7%", "93.7%", "0.861"],
    ],
    subtitle="Default 0.5 threshold — PR-AUC is the primary metric for imbalanced data",
    col_widths=[Inches(3.5), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7)],
    highlight_row=1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Full Model Comparison
# ═══════════════════════════════════════════════════════════════════════════
add_table_slide(prs,
    "Full Model Comparison (Validation Set)",
    ["Model", "Val PR-AUC", "Val ROC-AUC", "F1 Thresh", "Val F1", "Val Recall", "Min Cost (€)"],
    [
        ["LR (balanced)",       "0.752", "0.986", "0.990", "0.727", "86.5%", "5,380"],
        ["RF (balanced)",       "0.864", "0.951", "0.192", "0.877", "86.5%", "5,080"],
        ["GBT (balanced)",      "0.784", "0.982", "0.951", "0.853", "82.4%", "6,580"],
        ["GBT (SMOTE)",         "0.853", "0.973", "0.960", "0.879", "83.8%", "6,050"],
        ["DT (balanced)",       "0.656", "0.945", "0.990", "0.727", "86.5%", "5,380"],
        ["Isolation Forest",    "0.088", "0.945", "N/A",   "0.177", "20.3%", "N/A"],
    ],
    subtitle="★ Best overall: Random Forest (balanced) — highest PR-AUC + lowest business cost",
    col_widths=[Inches(2.3), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.3), Inches(1.6), Inches(1.7)],
    highlight_row=1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Threshold Tuning
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text(slide, "Dynamic Threshold Tuning", Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85),
         font_size=28, bold=True, color=WHITE)
add_rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)
add_rect(slide, 0, Inches(1.16), W, H - Inches(1.16), LIGHT_GRAY)

add_text(slide,
    "Why tune the threshold?",
    Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.4),
    font_size=18, bold=True, color=DARK_BLUE)
add_text(slide,
    "Sklearn classifiers default to 0.5, but the optimal decision boundary for imbalanced fraud data is much lower.\n"
    "We sweep threshold t ∈ [0, 1] on the validation set and pick the value that maximises F1 or minimises cost.",
    Inches(0.4), Inches(1.75), Inches(12.5), Inches(0.9),
    font_size=13, color=DARK_GRAY)

# 3-column explanation
cols = [
    ("F1-Optimal Threshold", DARK_BLUE,
     ["Sweep t on validation set",
      "Compute F1 at each t",
      "Pick t* = argmax F1",
      "RF best: t* = 0.192",
      "Lower threshold → more recalls"]),
    ("Cost-Optimal Threshold", MED_BLUE,
     ["Compute cost = €500×FN + €10×FP",
      "Pick t* = argmin cost",
      "RF best: t* = 0.044",
      "Very aggressive — catches more fraud",
      "Trades precision for coverage"]),
    ("Key Insight", GREEN,
     ["Default 0.5 misses ~20% of fraud",
      "RF F1-optimal (t=0.192) catches 86.5%",
      "Even at t=0.192, precision stays at 88.9%",
      "Threshold tuning is free (no retraining)",
      "Critical for real fraud systems"]),
]
col_w = Inches(4.0)
start_x = Inches(0.3)
for i, (title, color, pts) in enumerate(cols):
    x = start_x + i * (col_w + Inches(0.25))
    add_rect(slide, x, Inches(2.85), col_w, Inches(4.2), color)
    add_text(slide, title, x + Inches(0.1), Inches(2.95), col_w - Inches(0.2), Inches(0.45),
             font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    t = Inches(3.5)
    for p in pts:
        add_text(slide, "▸  " + p, x + Inches(0.15), t, col_w - Inches(0.3), Inches(0.52),
                 font_size=12, color=WHITE)
        t += Inches(0.55)

add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Cost-Sensitive Evaluation
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text(slide, "Cost-Sensitive Evaluation", Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85),
         font_size=28, bold=True, color=WHITE)
add_rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)
add_rect(slide, 0, Inches(1.16), W, H - Inches(1.16), LIGHT_GRAY)

# Cost formula box
add_rect(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(1.05), DARK_BLUE)
add_text(slide,
    "Business Cost  =  €500 × False Negatives  +  €10 × False Positives",
    Inches(0.4), Inches(1.42), Inches(12.5), Inches(0.55),
    font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(slide,
    "FN: missed fraud (severe)     |     FP: legitimate flagged (minor inconvenience)",
    Inches(0.4), Inches(1.95), Inches(12.5), Inches(0.35),
    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

# Validation cost table
headers = ["Model", "F1 Thresh", "Min Val Cost (€)", "Cost Thresh"]
rows = [
    ["LR (balanced)",  "0.990", "€5,380",  "0.985"],
    ["RF (balanced)",  "0.192", "€5,080 ★", "0.044"],
    ["GBT (balanced)", "0.951", "€6,580",  "0.837"],
    ["GBT (SMOTE)",    "0.960", "€6,050",  "0.251"],
    ["DT (balanced)",  "0.990", "€5,380",  "0.951"],
]
col_widths = [Inches(3.0), Inches(2.0), Inches(3.5), Inches(2.0)]
top = Inches(2.5)
left_start = Inches(1.7)
n_cols = len(headers)

x = left_start
for h in headers:
    add_rect(slide, x, top, col_widths[headers.index(h)], Inches(0.38), DARK_BLUE)
    add_text(slide, h, x + Inches(0.05), top + Inches(0.03),
             col_widths[headers.index(h)] - Inches(0.1), Inches(0.32),
             font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    x += col_widths[headers.index(h)]
top += Inches(0.38)

for ri, row in enumerate(rows):
    bg = RGBColor(0xE8, 0xF4, 0xFD) if ri == 1 else (WHITE if ri % 2 == 0 else LIGHT_GRAY)
    x = left_start
    for ci, cell in enumerate(row):
        add_rect(slide, x, top, col_widths[ci], Inches(0.38), bg,
                 line_color=RGBColor(0xCC, 0xCC, 0xCC))
        is_bold = ri == 1
        add_text(slide, cell, x + Inches(0.05), top + Inches(0.03),
                 col_widths[ci] - Inches(0.1), Inches(0.32),
                 font_size=11, bold=is_bold,
                 color=GREEN if ri == 1 else DARK_GRAY, align=PP_ALIGN.CENTER)
        x += col_widths[ci]
    top += Inches(0.38)

# Insight box
add_rect(slide, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.4), MED_BLUE)
add_text(slide,
    "Key Takeaway: RF (balanced) achieves the lowest validation cost of €5,080 at cost-optimal threshold t=0.044.\n"
    "The cost model forces us to prefer models that prioritize recall, even at the expense of slightly lower precision.",
    Inches(0.5), Inches(5.72), Inches(12.3), Inches(1.0),
    font_size=14, color=WHITE)

add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — OOD Detection (Novel Contribution)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text(slide, "Novel Contribution: OOD Detection with Isolation Forest",
         Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85),
         font_size=26, bold=True, color=WHITE)
add_rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)
add_rect(slide, 0, Inches(1.16), W, H - Inches(1.16), LIGHT_GRAY)

# Left: Motivation
add_rect(slide, Inches(0.3), Inches(1.3), Inches(5.8), Inches(5.7), DARK_BLUE)
add_text(slide, "Motivation", Inches(0.45), Inches(1.42), Inches(5.5), Inches(0.4),
         font_size=16, bold=True, color=GOLD)
pts_l = [
    "Supervised models are trained on known fraud patterns",
    "Novel fraud tactics not in training data will be missed",
    "Isolation Forest is trained ONLY on legitimate transactions",
    "Any transaction that looks 'anomalous' relative to legitimate data is flagged",
    "Complements supervised models — catches what they cannot",
    "Unsupervised: requires no fraud labels",
]
t = Inches(1.95)
for p in pts_l:
    add_text(slide, "▸  " + p, Inches(0.45), t, Inches(5.5), Inches(0.6),
             font_size=12, color=WHITE)
    t += Inches(0.62)

# Right: How it works + results
add_rect(slide, Inches(6.4), Inches(1.3), Inches(6.6), Inches(2.6), MED_BLUE)
add_text(slide, "How Isolation Forest Works", Inches(6.55), Inches(1.42), Inches(6.3), Inches(0.4),
         font_size=15, bold=True, color=GOLD)
how = [
    "Builds random trees that recursively partition features",
    "Anomalies are isolated in fewer splits (shorter paths)",
    "Anomaly score = average path length across trees",
    "contamination=0.01 (flags top 1% as anomalous)",
]
t = Inches(1.95)
for p in how:
    add_text(slide, "•  " + p, Inches(6.55), t, Inches(6.2), Inches(0.5),
             font_size=12, color=WHITE)
    t += Inches(0.52)

add_rect(slide, Inches(6.4), Inches(4.1), Inches(6.6), Inches(2.9), WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(slide, "OOD Detection Results", Inches(6.55), Inches(4.2), Inches(6.2), Inches(0.4),
         font_size=15, bold=True, color=DARK_BLUE)
ood_res = [
    ("ROC-AUC",            "0.945   (competitive!)"),
    ("PR-AUC",             "0.088   (low — no fraud in training)"),
    ("Recall (fraud)",     "20.3%   (unsupervised baseline)"),
    ("Interpretation",     "Used as a parallel detector, not a replacement"),
    ("Best use case",      "Flag transactions that evade supervised models"),
]
t = Inches(4.75)
for label, val in ood_res:
    add_text(slide, label + ":", Inches(6.6), t, Inches(2.5), Inches(0.4),
             font_size=11, bold=True, color=DARK_BLUE)
    add_text(slide, val, Inches(9.2), t, Inches(3.6), Inches(0.4),
             font_size=11, color=DARK_GRAY)
    t += Inches(0.44)

add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Error Analysis
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide(prs, "Error Analysis — Where the Model Fails",
    [
        "False Negatives (Missed Fraud) — most critical errors",
        "  Fraudulent transactions with Amount ≈ legitimate median → hard to distinguish by value alone",
        "  Fraud samples near the class boundary in PCA feature space",
        "  Rare fraud subtypes not well-represented in training data",
        "False Positives (Legitimate Flagged) — secondary concern",
        "  Legitimate high-value transactions that look statistically unusual",
        "  Transactions at unusual hours (low Time feature values)",
        "  Outlier legitimate users with non-standard spending patterns",
        "Error Patterns by Model",
        "  LR: High recall but poor precision at default threshold — floods alerts",
        "  RF: Best balance — low FP rate at t=0.192 with strong recall",
        "  DT: Overfits training; high variance in error distribution",
    ],
    subtitle="Analysis of False Negatives and False Positives from the best model (RF balanced)")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Final Test Results
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text(slide, "Final Test Set Results — Best Model: RF (balanced)",
         Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85),
         font_size=27, bold=True, color=WHITE)
add_rect(slide, 0, Inches(1.1), W, Inches(0.06), GOLD)
add_rect(slide, 0, Inches(1.16), W, H - Inches(1.16), LIGHT_GRAY)

add_text(slide,
    "Test set used exactly ONCE — these are the final, unbiased results",
    Inches(0.4), Inches(1.28), Inches(12.5), Inches(0.35),
    font_size=14, italic=True, color=RED, align=PP_ALIGN.CENTER)

# Metric cards
metrics_cards = [
    ("PR-AUC",        "0.8116",  "Primary metric"),
    ("ROC-AUC",       "0.9376",  "Discrimination"),
    ("F1-Score",      "0.8356",  "Balanced measure"),
    ("Recall",        "82.4%",   "Fraud caught"),
    ("Precision",     "84.7%",   "Alert accuracy"),
    ("Business Cost", "€6,610",  "Real-world cost"),
]
card_w = Inches(2.0)
card_h = Inches(1.55)
start_x = Inches(0.35)
t = Inches(1.75)

for i, (label, value, sub) in enumerate(metrics_cards):
    x = start_x + i * (card_w + Inches(0.18))
    color = GREEN if label in ("PR-AUC", "Recall", "F1-Score") else DARK_BLUE
    add_rect(slide, x, t, card_w, card_h, color)
    add_text(slide, value, x, t + Inches(0.12), card_w, Inches(0.7),
             font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, t + Inches(0.82), card_w, Inches(0.38),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sub, x, t + Inches(1.18), card_w, Inches(0.3),
             font_size=10, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER, italic=True)

# Confusion matrix summary
add_rect(slide, Inches(0.3), Inches(3.55), Inches(6.0), Inches(3.3), DARK_BLUE)
add_text(slide, "Confusion Matrix Summary (Test Set)", Inches(0.45), Inches(3.65),
         Inches(5.7), Inches(0.4), font_size=14, bold=True, color=GOLD)

cm_data = [
    ("True Negatives (TN)", "56,852", "Correctly flagged as legitimate"),
    ("False Positives (FP)", "17", "Legitimate flagged as fraud"),
    ("False Negatives (FN)", "17", "Missed fraud transactions"),
    ("True Positives (TP)", "80", "Correctly caught fraud"),
]
t2 = Inches(4.2)
for label, count, desc in cm_data:
    col = GREEN if label.startswith("True") else RED
    add_rect(slide, Inches(0.45), t2, Inches(0.6), Inches(0.52), col)
    add_text(slide, count, Inches(0.45), t2, Inches(0.6), Inches(0.52),
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, Inches(1.15), t2 + Inches(0.02), Inches(2.5), Inches(0.26),
             font_size=11, bold=True, color=WHITE)
    add_text(slide, desc, Inches(1.15), t2 + Inches(0.28), Inches(4.7), Inches(0.22),
             font_size=10, color=RGBColor(0xBB, 0xCC, 0xDD))
    t2 += Inches(0.6)

# Observations
add_rect(slide, Inches(6.6), Inches(3.55), Inches(6.4), Inches(3.3), WHITE,
         line_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(slide, "Key Observations", Inches(6.8), Inches(3.65), Inches(6.0), Inches(0.4),
         font_size=14, bold=True, color=DARK_BLUE)
obs = [
    "Caught 80 out of 97 fraud cases in test set",
    "Only 17 missed frauds (FN) at threshold t=0.192",
    "Only 17 false alarms (FP) out of 56,869 legitimate",
    "Val → Test generalization gap is small (~3% PR-AUC)",
    "Model shows no sign of severe overfitting",
    "Cost of €6,610 is acceptable for 56K transactions",
]
t2 = Inches(4.2)
for o in obs:
    add_text(slide, "▸  " + o, Inches(6.8), t2, Inches(6.0), Inches(0.46),
             font_size=12, color=DARK_GRAY)
    t2 += Inches(0.5)

add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusions
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, DARK_BLUE)
add_rect(slide, 0, Inches(1.1), W, Inches(0.08), GOLD)
add_rect(slide, 0, H - Inches(1.3), W, Inches(0.08), GOLD)

add_text(slide, "Conclusions & Future Work",
         Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8),
         font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Two columns
for col_x, title, color, items in [
    (Inches(0.3), "Key Findings", GOLD, [
        "RF (balanced) is the best model: PR-AUC 0.864, F1 0.877",
        "Dynamic thresholding (t=0.192) is critical — default 0.5 under-performs",
        "Balanced class weights outperform SMOTE and under-sampling",
        "Cost-sensitive threshold (t=0.044) further reduces business cost to €5,080",
        "Isolation Forest provides a useful unsupervised safety net",
        "Val → Test gap is small, indicating good generalization",
    ]),
    (Inches(6.9), "Future Work", GOLD, [
        "Ensemble: combine RF + Isolation Forest scores",
        "Online learning: retrain on streaming transactions",
        "Explore XGBoost / LightGBM for further PR-AUC gains",
        "Calibrate predicted probabilities (Platt scaling)",
        "Use SHAP values for transaction-level explainability",
        "Deploy as a real-time scoring microservice",
    ]),
]:
    add_rect(slide, col_x, Inches(1.3), Inches(6.2), Inches(5.7), MED_BLUE)
    add_text(slide, title, col_x + Inches(0.15), Inches(1.42), Inches(5.9), Inches(0.45),
             font_size=17, bold=True, color=GOLD)
    t = Inches(2.0)
    for item in items:
        add_text(slide, "▸  " + item, col_x + Inches(0.2), t, Inches(5.8), Inches(0.62),
                 font_size=12, color=WHITE)
        t += Inches(0.65)

add_text(slide,
    "Thank You  |  Questions?",
    Inches(0.5), H - Inches(1.2), Inches(12.3), Inches(0.6),
    font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(slide,
    "Om Suresh Prajapati  ·  oprajapa@usc.edu  ·  EE559 Spring 2026",
    Inches(0.5), H - Inches(0.75), Inches(12.3), Inches(0.4),
    font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/omsureshprajapati/ML-1/Project/reports/EE559_Fraud_Detection_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
